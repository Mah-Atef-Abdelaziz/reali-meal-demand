import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_deck():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette Definitions
    BG_DARK = RGBColor(18, 19, 24)       # Charcoal Dark #121318
    CARD_BG = RGBColor(26, 28, 35)       # Card BG #1A1C23
    GOLD = RGBColor(222, 193, 92)        # Gold #DEC15C
    WHITE = RGBColor(255, 255, 255)      # White
    MUTED = RGBColor(163, 163, 163)      # Light Grey #A3A3A3
    BORDER_COLOR = RGBColor(60, 62, 74)  # Border Grey

    def apply_slide_bg(slide):
        # Add full-screen rectangle for background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background() # No line
        return bg

    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Montserrat'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = GOLD
        return title_box

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide1)

    # Accent Graphic (Gold block)
    accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0

    p1 = tf.paragraphs[0]
    p1.text = "REAL.i"
    p1.font.name = 'Montserrat'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = GOLD

    p2 = tf.add_paragraph()
    p2.text = "MEAL DEMAND AI FORECASTING PLATFORM"
    p2.font.name = 'Montserrat'
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "Enterprise AI System to Optimize Catering Operations, Reduce Waste & Lower Costs"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(15)

    p4 = tf.add_paragraph()
    p4.text = "Saudi Arabia • 15 Corporate Locations • 1.47M Transaction Records"
    p4.font.name = 'Calibri'
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = GOLD
    p4.space_before = Pt(30)


    # -------------------------------------------------------------
    # SLIDE 2: The Core Problem
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide2)
    add_title(slide2, "The Catering Scale Challenge")

    # Left Card: The Problem
    card1 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = BORDER_COLOR
    
    tb1 = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_bottom = tf1.margin_right = 0
    
    p = tf1.paragraphs[0]
    p.text = "Operational Difficulties"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    bullets1 = [
        "Traditional catering relies on static headcounts, leading to massive inefficiencies.",
        "Over-preparation results in a baseline 30% food waste rate.",
        "Under-preparation leads to portion shortages and damages employee satisfaction.",
        "Industrial plants & offshore rigs suffer from complex 14-day shift rotation schedules."
    ]
    for b in bullets1:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Calibri'
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # Right Card: The Variables
    card2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = BORDER_COLOR

    tb2 = slide2.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_bottom = tf2.margin_right = 0

    p = tf2.paragraphs[0]
    p.text = "Key Forecasting Drivers"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    bullets2 = [
        "Rotational Shifts: Day/Night rotations dynamically impact headcount.",
        "Weather Conditions: Hot temperatures shift demand from hot food to salads.",
        "Calendar Events: Saudi national holidays & corporate workshops.",
        "Location Capacities: Distinguishing office space from remote industrial sites."
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Calibri'
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)


    # -------------------------------------------------------------
    # SLIDE 3: System Architecture
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide3)
    add_title(slide3, "Decoupled System Architecture")

    # Draw 3 column boxes (Frontend, Backend, Database)
    col_width = Inches(3.64)
    gap = Inches(0.4)
    start_x = Inches(0.8)

    layers = [
        ("1. Frontend UI", "Next.js 16 Console", [
            "React 19 App Router dashboard.",
            "Gold / charcoal micro-animations.",
            "Fully mobile-responsive sidebar with '<' and '^' collapsible triggers.",
            "Displays real-time KPIs, savings charts, and predictions."
        ]),
        ("2. Backend API", "FastAPI Core Engine", [
            "Async runtime using Python 3.12.",
            "Asyncpg pipeline for non-blocking database queries.",
            "OpenAI GPT-4o-mini LangChain RAG integrations.",
            "Hosted in Docker container on Hugging Face Spaces (Port 7860)."
        ]),
        ("3. Database Layer", "Neon Serverless Postgres", [
            "Fully normalized 18-table 3NF relational database.",
            "Contains 1.47 million historical records.",
            "Handles audit logs, predictions, weather, and employee rosters.",
            "Auto-scaling database keeps infrastructure cost at zero."
        ])
    ]

    for idx, (title, sub, bullets) in enumerate(layers):
        x = start_x + idx * (col_width + gap)
        # Background card
        card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), col_width, Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tb = slide3.shapes.add_textbox(x + Inches(0.25), Inches(2.0), col_width - Inches(0.5), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Montserrat'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GOLD

        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.name = 'Calibri'
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(4)

        for b in bullets:
            p = tf.add_paragraph()
            p.text = "• " + b
            p.font.name = 'Calibri'
            p.font.size = Pt(13)
            p.font.color.rgb = MUTED
            p.space_before = Pt(10)


    # -------------------------------------------------------------
    # SLIDE 4: Database Design (3NF Schema)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide4)
    add_title(slide4, "Normalized Database Design (3NF)")

    # Left: Database Summary Card
    card_left = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = BORDER_COLOR

    tb_left = slide4.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0

    p = tf_left.paragraphs[0]
    p.text = "Relational Integrity"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    db_points = [
        "18 normalized tables structured to avoid redundancy.",
        "Maintains full referential integrity constraints across models.",
        "Seeded with over 1.47 million production rows.",
        "Optimized indexes on transactional dates & location IDs."
    ]
    for pt in db_points:
        p = tf_left.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # Right: Schema Dictionary Card
    card_right = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = BORDER_COLOR

    tb_right = slide4.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_top = tf_right.margin_bottom = tf_right.margin_right = 0

    p = tf_right.paragraphs[0]
    p.text = "Key Database Tables"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    tables_list = [
        "meal_transactions: Individual meal logs (~1.03M rows)",
        "employees: Demographic profiles and rosters (~100k rows)",
        "attendance: Clock-in/out registers (~210k rows)",
        "daily_menus & menu_items: Menu scheduling records",
        "weather_data & holiday_calendar: External context tables",
        "prediction_results & model_logs: Machine learning logs"
    ]
    for tbl in tables_list:
        p = tf_right.add_paragraph()
        p.text = "• " + tbl
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.space_before = Pt(10)


    # -------------------------------------------------------------
    # SLIDE 5: Machine Learning & Preprocessing
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide5)
    add_title(slide5, "Machine Learning Preprocessing")

    # Left Card: Feature Engineering
    card_left = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = BORDER_COLOR

    tb_left = slide5.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0

    p = tf_left.paragraphs[0]
    p.text = "Feature Engineering (38 Vectors)"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    feat_points = [
        "Temporal vectors: Month, day of week, day of year, quarter.",
        "Saudi Work Week Adjustment: Custom 'saudi_dow' feature maps Sunday-Thursday work cycle correctly.",
        "Demand Lag Indicators: Lags of 1-day, 7-days, 14-days, and 28-days feed historical trends to model.",
        "Moving Averages: 7-day and 14-day rolling demand averages."
    ]
    for pt in feat_points:
        p = tf_left.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # Right Card: Context Vectors
    card_right = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = BORDER_COLOR

    tb_right = slide5.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_top = tf_right.margin_bottom = tf_right.margin_right = 0

    p = tf_right.paragraphs[0]
    p.text = "Environmental & External Inputs"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    context_points = [
        "Site Capacity: Max seating thresholds per location.",
        "Weather Vectors: Average temperature, humidity, and rainfall.",
        "Holiday Calendars: National holidays and company events.",
        "Visitor Log: Incorporating scheduled external visitors.",
        "Department Metrics: Roster counts mapped per department."
    ]
    for pt in context_points:
        p = tf_right.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.space_before = Pt(10)


    # -------------------------------------------------------------
    # SLIDE 6: Model Selection & Performance
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide6)
    add_title(slide6, "Model Performance & Accuracy")

    # Left: Evaluation Metrics Table
    card_left = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = BORDER_COLOR

    tb_left = slide6.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.9), Inches(4.2))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0

    p = tf_left.paragraphs[0]
    p.text = "Forecasting Algorithm Comparison"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p_tbl = tf_left.add_paragraph()
    p_tbl.text = "1. XGBoost Regressor (Selected)\n    R² Score: 0.9820  |  MAE: 3.75 meals  |  MAPE: 22.21%\n\n" + \
                 "2. Random Forest Regressor\n    R² Score: 0.9724  |  MAE: 5.12 meals  |  MAPE: 28.54%\n\n" + \
                 "3. Gradient Boosting Regressor\n    R² Score: 0.9680  |  MAE: 5.95 meals  |  MAPE: 31.02%\n\n" + \
                 "4. Linear Regression Baseline\n    R² Score: 0.9148  |  MAE: 9.49 meals  |  MAPE: 44.80%"
    p_tbl.font.name = 'Calibri'
    p_tbl.font.size = Pt(14)
    p_tbl.font.color.rgb = WHITE
    p_tbl.space_before = Pt(15)

    # Right: Summary Box
    card_right = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(1.8), Inches(4.8), Inches(4.8))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = GOLD # Highlight border
    
    tb_right = slide6.shapes.add_textbox(Inches(8.0), Inches(2.1), Inches(4.2), Inches(4.2))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_top = tf_right.margin_bottom = tf_right.margin_right = 0

    p = tf_right.paragraphs[0]
    p.text = "Why XGBoost Wins"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    why_points = [
        "Effectively captures complex non-linear relations and interaction variables.",
        "Highly robust against missing records or database outliers.",
        "Selected configuration: 500 estimators, max depth of 7, learning rate of 0.05.",
        "Model retrains automatically on nightly feedback loops."
    ]
    for pt in why_points:
        p = tf_right.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)


    # -------------------------------------------------------------
    # SLIDE 7: Conversational AI & RAG Engine
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide7)
    add_title(slide7, "Conversational AI & RAG Engine")

    # Left: AI Architecture
    card_left = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = BORDER_COLOR

    tb_left = slide7.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0

    p = tf_left.paragraphs[0]
    p.text = "Natural Language Interface"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    rag_points = [
        "Powered by LangChain orchestrator mapping to ChatOpenAI.",
        "Dynamically translates user queries into SQL queries.",
        "Executes queries on PostgreSQL db to fetch context statistics.",
        "Passes context to GPT-4o-mini for conversational answers.",
        "Features local Llama-3 fallback path if API key is missing."
    ]
    for pt in rag_points:
        p = tf_left.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # Right: Examples Box
    card_right = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = BORDER_COLOR

    tb_right = slide7.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_top = tf_right.margin_bottom = tf_right.margin_right = 0

    p = tf_right.paragraphs[0]
    p.text = "Example Assistant Interactions"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    examples = [
        "Q: 'How many lunch meals do we need for ONS-JBL tomorrow?'",
        "A: 'Based on active rosters, shift structures, and tomorrow's warm weather forecast, the model predicts a demand of 1,320 lunch meals.'",
        "",
        "Q: 'Show me our cost savings this month.'",
        "A: 'Our monthly cost savings has reached 148,500 SAR, reflecting a waste reduction of 24.5% compared to the baseline.'"
    ]
    for ex in examples:
        p = tf_right.add_paragraph()
        p.text = ex
        p.font.name = 'Calibri'
        p.font.size = Pt(13)
        if ex.startswith("Q:"):
            p.font.color.rgb = GOLD
            p.font.bold = True
            p.space_before = Pt(12)
        else:
            p.font.color.rgb = MUTED
            p.space_before = Pt(4)


    # -------------------------------------------------------------
    # SLIDE 8: Live System Console & Responsiveness
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide8)
    add_title(slide8, "Interactive Console & Mobile Design")

    # Left: Web UI Features
    card_left = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = BORDER_COLOR

    tb_left = slide8.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = 0

    p = tf_left.paragraphs[0]
    p.text = "Central Operations Console"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    console_points = [
        "Premium dashboard featuring gold/charcoal aesthetics.",
        "Interactive forecast visualization charts (Recharts).",
        "Live network connection state indicator (ONLINE vs DEMO).",
        "Menu management modules for active scheduling."
    ]
    for pt in console_points:
        p = tf_left.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # Right: Mobile Responsive Features
    card_right = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = GOLD # Highlight border

    tb_right = slide8.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_top = tf_right.margin_bottom = tf_right.margin_right = 0

    p = tf_right.paragraphs[0]
    p.text = "Mobile Layout Adaptability"
    p.font.name = 'Montserrat'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    mobile_points = [
        "Full layout adaptability to tablet & mobile devices.",
        "Desktop collapsing via '<' / '>' controls reduces the sidebar width to a compact icon-only view.",
        "Mobile collapsing via '^' overlay closes the sidebar vertically on phone viewports.",
        "Guarantees that menu navigation doesn't block analytics visibility on smaller screens."
    ]
    for pt in mobile_points:
        p = tf_right.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)


    # -------------------------------------------------------------
    # SLIDE 9: ROI & Commercial Sustainability
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_bg(slide9)
    add_title(slide9, "Commercial ROI & Sustainability")

    # Big Metrics Grid (3 columns)
    col_w = Inches(3.64)
    start_x = Inches(0.8)
    gap = Inches(0.4)

    metrics = [
        ("24.5%", "FOOD WASTE REDUCTION", "Average waste rate drops from 30.0% to 5.5% using XGBoost predictions."),
        ("148.5k SAR", "MONTHLY COST SAVINGS", "Based on average meal cost of 15 SAR across 15 operational locations."),
        ("44.5 tons", "CO2 FOOTPRINT OFFSET", "Monthly greenhouse gas emissions avoided by reducing daily food disposal.")
    ]

    for idx, (val, title, desc) in enumerate(metrics):
        x = start_x + idx * (col_w + gap)
        # Background card
        card = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), col_w, Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = GOLD if idx == 1 else BORDER_COLOR # Highlight savings card

        tb = slide9.shapes.add_textbox(x + Inches(0.25), Inches(2.2), col_w - Inches(0.5), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0

        # Big Number Value
        p = tf.paragraphs[0]
        p.text = val
        p.font.name = 'Montserrat'
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = GOLD

        # Title Label
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = 'Montserrat'
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER
        p_title.font.color.rgb = WHITE
        p_title.space_before = Pt(15)

        # Description
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Calibri'
        p_desc.font.size = Pt(13)
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.font.color.rgb = MUTED
        p_desc.space_before = Pt(15)


    # Save Deck
    output_path = "docs/REAL.i_Meal_Demand_AI_Presentation.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    create_deck()
